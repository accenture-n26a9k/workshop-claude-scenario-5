"""
retheme.py — PPTX theme migration script

Replaces color tokens and slide masters with a new theme
while preserving existing PPTX slide content.

Usage (bash):
    VENV="$HOME/.claude/skills/.venv/Scripts/python.exe"

    # ACN -> Fiori
    $VENV ~/.claude/skills/acnpptx/scripts/retheme.py deck.pptx fiori

    # Fiori -> ACN (explicit source theme)
    $VENV ~/.claude/skills/acnpptx/scripts/retheme.py deck.pptx accenture --from fiori --out deck_acn.pptx

Python API:
    from retheme import migrate
    migrate("deck.pptx", target_theme="fiori", source_theme="accenture")

Processing steps:
  1. Build color token mapping (old hex -> new hex)
  2. Transplant slides onto the new template base
  3. Resolve schemeClr via source theme to srgbClr, then replace colors
  4. Transplant placeholders as free-floating textboxes if not present in new template
  5. Transplant image relationships (icon PNGs, etc.)

Limitations:
  - Colors inside charts are not replaced (charts are preserved as content but colors remain unchanged)
  - Decorations provided by slide master (logo, GT symbol, footer) are updated to new master
"""

import sys
import os
import copy
import json
import argparse
import glob
import colorsys

# Windows terminal encoding
if sys.platform == "win32":
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

_SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))
_SKILL_DIR   = os.path.dirname(_SCRIPTS_DIR)
_THEMES_DIR  = os.path.join(_SKILL_DIR, "assets", "themes")
_MASTERS_DIR = os.path.join(_SKILL_DIR, "assets", "masters")

sys.path.insert(0, _SCRIPTS_DIR)

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree as _et

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS   = {"a": _NS_A}
_PH_SKIP = {20, 21}   # footer / page-num — auto-provided by new master
_SHAPE_TAGS = frozenset([
    qn("p:sp"), qn("p:pic"), qn("p:graphicFrame"),
    qn("p:grpSp"), qn("p:cxnSp"),
])


# ── Utilities ─────────────────────────────────────────────────────────────

def _load_theme_data(theme_name: str) -> dict:
    path = os.path.join(_THEMES_DIR, f"{theme_name}.json")
    if not os.path.isfile(path):
        # case-insensitive search by name
        for f in glob.glob(os.path.join(_THEMES_DIR, "*.json")):
            with open(f, encoding="utf-8") as fp:
                d = json.load(fp)
            if d.get("name", "").lower() == theme_name.lower():
                return d
        raise FileNotFoundError(f"Theme not found: {theme_name} (looked in {_THEMES_DIR})")
    with open(path, encoding="utf-8") as f:
        return json.load(f)


def _resolve_template_path(theme_data: dict) -> str:
    tpl = theme_data.get("template", "")
    if tpl:
        if not os.path.isabs(tpl):
            tpl = os.path.join(_SKILL_DIR, tpl)
        if os.path.isfile(tpl):
            return tpl
    # fallback: masters/<name>.pptx
    name = theme_data.get("name", "").lower()
    tpl = os.path.join(_MASTERS_DIR, f"{name}.pptx")
    if os.path.isfile(tpl):
        return tpl
    raise FileNotFoundError(
        f"Template PPTX not found for theme '{theme_data.get('name')}'. "
        f"Expected: {tpl}"
    )


def _build_color_map(src_tokens: dict, dst_tokens: dict) -> dict:
    """Build old-hex -> new-hex map via token names."""
    cmap = {}
    for token, src_hex in src_tokens.items():
        if token in dst_tokens:
            s = src_hex.lstrip("#").upper()
            d = dst_tokens[token].lstrip("#").upper()
            if s != d:
                cmap[s] = d
    return cmap


def _expand_token_colors(tokens: dict) -> dict:
    """
    Compute intermediate colors from tokens and return an expanded {alias: hex(6)} map.
    Uses the same calculation as theme_selector.py _mid / _mid2
    to generate midpoint colors equivalent to LIGHT_PURPLE/DARK_PURPLE.
    """
    def _mid(h1, h2):
        h1, h2 = h1.lstrip("#"), h2.lstrip("#")
        r = (int(h1[0:2], 16) + int(h2[0:2], 16)) // 2
        g = (int(h1[2:4], 16) + int(h2[2:4], 16)) // 2
        b = (int(h1[4:6], 16) + int(h2[4:6], 16)) // 2
        return f"{r:02X}{g:02X}{b:02X}"

    expanded = {k: v.lstrip("#").upper() for k, v in tokens.items()}
    if "primary" in tokens and "primary_dark" in tokens:
        expanded["_mid_pd"] = _mid(tokens["primary"], tokens["primary_dark"])
    if "primary" in tokens and "primary_light" in tokens:
        expanded["_mid_pl"] = _mid(tokens["primary"], tokens["primary_light"])
    if "primary_dark" in tokens and "primary_light" in tokens:
        expanded["_mid_dl"] = _mid(tokens["primary_dark"], tokens["primary_light"])
    return expanded


def _collect_slide_colors(src_prs) -> set:
    """Return the set of actually used srgbClr hex values from all slides in the source PPTX."""
    colors = set()
    for slide in src_prs.slides:
        for el in slide._element.iter(f"{{{_NS_A}}}srgbClr"):
            val = el.get("val", "").upper()
            if len(val) == 6:
                colors.add(val)
    return colors


def _saturation(hex_color: str) -> float:
    """Return HSL saturation (0.0-1.0) from RGB hex."""
    r = int(hex_color[0:2], 16) / 255
    g = int(hex_color[2:4], 16) / 255
    b = int(hex_color[4:6], 16) / 255
    _, _, s = colorsys.rgb_to_hls(r, g, b)
    return s


def _build_color_map_fuzzy(
    src_tokens: dict, dst_tokens: dict,
    actual_colors: set, delta: int = 25,
    direct_delta: int = 442,
) -> dict:
    """
    Build color_map in 3 phases:

      1. Token name exact match (primary -> primary, etc.)
      2. Fuzzy match including src intermediate colors (when actual slide colors are close to src_tokens midpoints)
      3. Generic direct match: convert chromatic actual_colors to nearest dst_tokens chromatic color
         A generic approach that works even when src_tokens are unknown or outdated.
         Neutral colors (saturation < 0.15) are excluded.
         The default direct_delta is the maximum RGB distance (~442), so
         chromatic colors are always converted to the nearest dst chromatic color.

    Args:
        src_tokens:   Source theme tokens dict (may be empty)
        dst_tokens:   Target theme tokens dict
        actual_colors: Actual srgbClr hex set from slides (uppercase 6-digit)
        delta:        Max RGB Euclidean distance for Phase 2 (default 25)
        direct_delta: Max RGB Euclidean distance for Phase 3 (default 442 = unlimited)
    """
    _SAT_THRESHOLD = 0.15  # Below this saturation threshold, treat as neutral (black/white/gray) and skip conversion

    def _dist(h1: str, h2: str) -> float:
        r1, g1, b1 = int(h1[0:2], 16), int(h1[2:4], 16), int(h1[4:6], 16)
        r2, g2, b2 = int(h2[0:2], 16), int(h2[2:4], 16), int(h2[4:6], 16)
        return ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5

    # ── Phase 1: token name exact match ──────────────────────────────────────
    cmap = _build_color_map(src_tokens, dst_tokens)

    if not actual_colors:
        return cmap

    # ── Phase 2: fuzzy match including src intermediate colors ──────────────────────────────────
    if src_tokens:
        src_exp = _expand_token_colors(src_tokens)
        dst_exp = _expand_token_colors(dst_tokens)
        src_to_dst: dict[str, str] = {}
        for key, src_hex in src_exp.items():
            if key in dst_exp and src_hex != dst_exp[key]:
                src_to_dst[src_hex] = dst_exp[key]

        if src_to_dst:
            src_keys = list(src_to_dst.keys())
            fuzzy_added = 0
            for actual in actual_colors:
                if actual in cmap:
                    continue
                best_key, best_d = None, float("inf")
                for sk in src_keys:
                    d = _dist(actual, sk)
                    if d < best_d:
                        best_d, best_key = d, sk
                if best_key is not None and best_d <= delta:
                    cmap[actual] = src_to_dst[best_key]
                    fuzzy_added += 1
            if fuzzy_added:
                print(f"  Fuzzy color matches (src-based): {fuzzy_added} color(s)")

    # ── Phase 3: generic direct match (works with any source theme) ──────────────
    if dst_tokens:
        dst_exp = _expand_token_colors(dst_tokens)
        # dst chromatic color list (neutral colors excluded from conversion targets)
        dst_chromatic_vals = [v for v in dst_exp.values() if _saturation(v) >= _SAT_THRESHOLD]
        dst_hex_set = set(dst_chromatic_vals)  # Already a dst theme color, no conversion needed

        direct_added = 0
        for actual in actual_colors:
            if actual in cmap:
                continue  # Already a conversion target from Phase 1/2
            if actual in dst_hex_set:
                continue  # Already a dst theme color, no conversion needed
            if _saturation(actual) < _SAT_THRESHOLD:
                continue  # Skip neutral colors (black/white/gray)

            best_dst, best_d = None, float("inf")
            for dst_hex in dst_chromatic_vals:
                d = _dist(actual, dst_hex)
                if d < best_d:
                    best_d, best_dst = d, dst_hex
            if best_dst is not None and best_d <= direct_delta:
                cmap[actual] = best_dst
                direct_added += 1

        if direct_added:
            print(f"  Direct chromatic matches: {direct_added} color(s)")

    return cmap


def _replace_colors(element, color_map: dict) -> int:
    """Replace all srgbClr in the XML tree per color_map. Returns replacement count."""
    count = 0
    for el in element.iter(f"{{{_NS_A}}}srgbClr"):
        val = el.get("val", "").upper()
        if val in color_map:
            el.set("val", color_map[val])
            count += 1
    return count


# ── Layout type estimation ────────────────────────────────────────────────────────

_COVER_KEYWORDS   = {"cover", "title slide", "タイトル スライド", "タイトルスライド", "表紙"}
_SECTION_KEYWORDS = {"section", "セクション", "divider", "chapter", "区切り", "section header"}
_BLANK_KEYWORDS   = {"blank", "空白"}


def _classify_layout_type(layout_name: str) -> str:
    """Estimate slide type (cover/section/content/blank) from layout name."""
    name = layout_name.lower()
    for kw in _COVER_KEYWORDS:
        if kw in name:
            return "cover"
    for kw in _SECTION_KEYWORDS:
        if kw in name:
            return "section"
    for kw in _BLANK_KEYWORDS:
        if kw in name:
            return "blank"
    return "content"


def _build_dst_layout_map(dst_prs) -> dict:
    """
    Determine target template layouts by placeholder structure and
    map type -> layout. Works across different templates since it does not depend on layout names.

    Detection logic:
      cover   : Has both CENTER_TITLE (type=3) + SUBTITLE (type=4)
      section : Has CENTER_TITLE (type=3) but not SUBTITLE (type=4)
      content : Has TITLE (type=1) + BODY (type=2) + footer idx (20 or 21)
      blank   : Falls back to same layout as content
    """
    # placeholder type constants (pp_placeholder_type integer values)
    _CT  = 3   # CENTER_TITLE
    _T   = 1   # TITLE
    _ST  = 4   # SUBTITLE
    _B   = 2   # BODY

    layout_map: dict = {}

    for layout in dst_prs.slide_layouts:
        ph_types = {ph.placeholder_format.type for ph in layout.placeholders}
        ph_idxs  = {ph.placeholder_format.idx  for ph in layout.placeholders}

        # cover: CENTER_TITLE + SUBTITLE
        if "cover" not in layout_map:
            if _CT in ph_types and _ST in ph_types:
                layout_map["cover"] = layout

        # section: CENTER_TITLE, no SUBTITLE
        if "section" not in layout_map:
            if _CT in ph_types and _ST not in ph_types:
                layout_map["section"] = layout

        # content: TITLE (not CENTER_TITLE), BODY, footer idx=20 or 21
        if "content" not in layout_map:
            if (_T in ph_types and _CT not in ph_types and
                    _B in ph_types and (20 in ph_idxs or 21 in ph_idxs)):
                layout_map["content"] = layout

        # content_fallback1: TITLE + BODY (no footer)
        if "content" not in layout_map and "_content_fb1" not in layout_map:
            if _T in ph_types and _CT not in ph_types and _B in ph_types:
                layout_map["_content_fb1"] = layout

        # content_fallback2: TITLE only (simplest layout, equivalent to "Title Only")
        # Prefer layout with fewest placeholders (fewer decorative elements)
        if "content" not in layout_map:
            if (_T in ph_types and _CT not in ph_types and
                    _B not in ph_types and _ST not in ph_types):
                prev = layout_map.get("_content_fb2")
                if prev is None or len(list(layout.placeholders)) < len(list(prev.placeholders)):
                    layout_map["_content_fb2"] = layout

    # Staged fallback for content
    # fb2 (Title Only) is the cleanest canvas -> prioritize over fb1 (Title+Body)
    if "content" not in layout_map:
        if "_content_fb2" in layout_map:
            layout_map["content"] = layout_map["_content_fb2"]
        elif "_content_fb1" in layout_map:
            layout_map["content"] = layout_map["_content_fb1"]
    # Remove temporary keys
    layout_map.pop("_content_fb1", None)
    layout_map.pop("_content_fb2", None)

    # fallback by index
    n = len(dst_prs.slide_layouts)
    if "cover" not in layout_map:
        layout_map["cover"] = dst_prs.slide_layouts[0]
    if "content" not in layout_map:
        # Final fallback: idx=1 (safe for templates where "Title Only" is common)
        layout_map["content"] = dst_prs.slide_layouts[min(1, n - 1)]
    if "section" not in layout_map:
        layout_map["section"] = dst_prs.slide_layouts[min(6, n - 1)]
    layout_map["blank"] = layout_map.get("content", dst_prs.slide_layouts[0])

    print("Destination layout map:")
    for lt, layout in layout_map.items():
        print(f"  {lt:10s} → [{layout.name}]")
    print()

    return layout_map


def _remove_empty_placeholders(slide):
    """
    Remove empty text placeholders from slide XML.
    This prevents hint text like "Click to add title" from
    appearing in the rendered output.
    Target: placeholders with text frames but no actual text.
    Excluded: footer/pagenum (idx 20/21), those without text frames.
    """
    sp_tree = slide.shapes._spTree
    for el in list(sp_tree):
        if not _is_placeholder(el):
            continue
        idx = _get_ph_idx(el)
        if idx in _PH_SKIP:
            continue
        tx_body = el.find(f".//{qn('p:txBody')}")
        if tx_body is None:
            continue  # Skip image placeholders etc.
        all_text = "".join(
            (r.text or "") for r in tx_body.iter(f"{{{_NS_A}}}r")
        )
        if not all_text.strip():
            sp_tree.remove(el)


def _is_placeholder(el) -> bool:
    return el.find(f".//{qn('p:ph')}") is not None


def _get_ph_idx(el) -> int | None:
    """Return idx of p:ph element, or None if absent."""
    ph = el.find(f".//{qn('p:ph')}")
    if ph is None:
        return None
    return int(ph.get("idx", "0"))


# ── schemeClr resolution ─────────────────────────────────────────────────────────────

def _build_src_theme_map(src_prs) -> dict:
    """
    Build schemeClr slot name -> hex from source PPTX master theme XML.
    Same logic as master_to_theme._build_theme_map.
    """
    theme_map: dict = {}
    for master in src_prs.slide_masters:
        for rel in master.part.rels.values():
            if "theme" not in rel.reltype.lower():
                continue
            tree = _et.fromstring(rel.target_part.blob)
            clr_el = tree.find(".//a:clrScheme", _NS)
            if clr_el is None:
                continue
            for child in clr_el:
                tag = child.tag.split("}")[1]
                for sub in child:
                    val = sub.get("val") or sub.get("lastClr", "")
                    if not val:
                        continue
                    if val.upper() in ("WINDOWTEXT", "WINDOW"):
                        theme_map[tag] = "#000000" if "dk" in tag.lower() or tag == "tx1" else "#FFFFFF"
                    else:
                        theme_map[tag] = "#" + val.upper()
            break
    if "bg1" not in theme_map:
        theme_map["bg1"] = theme_map.get("lt1", "#FFFFFF")
    if "tx1" not in theme_map:
        theme_map["tx1"] = theme_map.get("dk1", "#000000")
    return theme_map


def _resolve_scheme_colors(element, src_theme_map: dict, color_map: dict):
    """
    Resolve all schemeClr in the XML tree to hex via src_theme_map, convert to srgbClr,
    and apply color_map. Targets not just solidFill but also gradFill/gs, tcPr, etc.
    """
    _FILL_CONTAINERS = {
        f"{{{_NS_A}}}solidFill",
        f"{{{_NS_A}}}gs",        # gradFill/gsLst/gs(gradient stop)
    }
    for sf in element.iter():
        if sf.tag not in _FILL_CONTAINERS:
            continue
        for child in list(sf):
            tag = child.tag.split("}")[1]
            if tag != "schemeClr":
                continue
            base = src_theme_map.get(child.get("val", ""))
            if not base:
                continue

            # Apply lumMod / lumOff / shade / tint in HLS color space
            r0 = int(base[1:3], 16)
            g0 = int(base[3:5], 16)
            b0 = int(base[5:7], 16)
            h, l, s = colorsys.rgb_to_hls(r0 / 255, g0 / 255, b0 / 255)
            for mod in child:
                mod_tag = mod.tag.split("}")[1]
                v = int(mod.get("val", "100000")) / 100000
                if mod_tag == "lumMod":
                    l = max(0.0, min(1.0, l * v))
                elif mod_tag == "lumOff":
                    l = max(0.0, min(1.0, l + v))
                elif mod_tag == "shade":
                    l = max(0.0, min(1.0, l * v))
                elif mod_tag == "tint":
                    l = max(0.0, min(1.0, l + (1 - l) * v))
            r1, g1, b1 = colorsys.hls_to_rgb(h, l, s)
            resolved = f"{int(r1 * 255):02X}{int(g1 * 255):02X}{int(b1 * 255):02X}"

            # Apply color_map replacement
            final = color_map.get(resolved, resolved)

            # Replace schemeClr with srgbClr
            sf.remove(child)
            srgb = _et.SubElement(sf, f"{{{_NS_A}}}srgbClr")
            srgb.set("val", final)


# ── Slide copy ─────────────────────────────────────────────────────────────

def _copy_image_rel(src_part, dst_part, old_rId: str) -> str | None:
    """Copy image from src_part to dst_part and return new rId."""
    try:
        rel = src_part.rels.get(old_rId)
        if rel is None:
            return None
        return dst_part.relate_to(rel.target_part, rel.reltype)
    except Exception as e:
        print(f"    Warning: image copy failed ({old_rId}): {e}")
        return None


def _fix_pic_rels(el_copy, src_slide, dst_slide):
    """
    Transplant r:embed rId for all p:pic elements in el_copy tree from src to dst.
    Handles all patterns: direct p:pic, inside groups, and post-placeholder conversion.
    """
    for pic in el_copy.iter(qn("p:pic")):
        blip = pic.find(f".//{{{_NS_A}}}blip")
        if blip is None:
            continue
        old_rId = blip.get(qn("r:embed"))
        if not old_rId:
            continue
        new_rId = _copy_image_rel(src_slide.part, dst_slide.part, old_rId)
        if new_rId:
            blip.set(qn("r:embed"), new_rId)


def _copy_shape_element(el, src_slide, dst_slide, dst_tree):
    """
    Add one shape element to dst_tree (including image/chart rId transplant).
    Process both direct p:pic and nested p:pic inside p:grpSp via _fix_pic_rels.
    """
    el_copy = copy.deepcopy(el)

    # Image (p:pic) / group shape (p:grpSp) — transplant r:embed for all p:pic
    if el.tag in (qn("p:pic"), qn("p:grpSp")):
        _fix_pic_rels(el_copy, src_slide, dst_slide)

    # graphicFrame (chart/table) — rId transplant
    elif el.tag == qn("p:graphicFrame"):
        for rel_el in el.iter():
            chart_rId = rel_el.get(qn("r:id"))
            if chart_rId and chart_rId.startswith("rId"):
                rel = src_slide.part.rels.get(chart_rId)
                if rel:
                    try:
                        new_rId = dst_slide.part.relate_to(rel.target_part, rel.reltype)
                        el_copy_rel_el = el_copy.find(
                            ".//" + rel_el.tag
                        ) if rel_el is not el else el_copy
                        if el_copy_rel_el is not None:
                            el_copy_rel_el.set(qn("r:id"), new_rId)
                    except Exception:
                        pass

    dst_tree.append(el_copy)


def _inject_xfrm(el_copy, sh):
    """
    Explicitly embed python-pptx shape resolved coordinates into el_copy p:spPr/a:xfrm.
    Needed when position is inherited from layout only and not recorded in slide XML.
    """
    if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
        return
    spPr = el_copy.find(qn("p:spPr"))
    if spPr is None:
        return
    # Remove existing xfrm
    existing = spPr.find(f"{{{_NS_A}}}xfrm")
    if existing is not None:
        spPr.remove(existing)
    # Insert new xfrm at beginning of spPr
    xfrm = _et.Element(f"{{{_NS_A}}}xfrm")
    off  = _et.SubElement(xfrm, f"{{{_NS_A}}}off")
    off.set("x", str(sh.left))
    off.set("y", str(sh.top))
    ext  = _et.SubElement(xfrm, f"{{{_NS_A}}}ext")
    ext.set("cx", str(sh.width))
    ext.set("cy", str(sh.height))
    spPr.insert(0, xfrm)


def _copy_shapes(src_slide, dst_slide, dst_ph_idxs: set):
    """
    Copy non-placeholder shapes from src to dst.
    Additionally, placeholders with idx not present in dst have p:ph removed and are
    transplanted as free-floating textboxes (prevents cover/message line disappearance).
    When coordinates are layout-inherited, explicitly embed python-pptx resolved values.
    """
    src_tree = src_slide.shapes._spTree
    dst_tree = dst_slide.shapes._spTree

    # Map python-pptx shapes by idx -> shape (for position resolution)
    src_ph_map = {
        sh.placeholder_format.idx: sh
        for sh in src_slide.shapes
        if sh.is_placeholder
    }

    # Remove default non-placeholder shapes from dst
    for el in list(dst_tree):
        if el.tag in _SHAPE_TAGS and not _is_placeholder(el):
            dst_tree.remove(el)

    for el in src_tree:
        if el.tag not in _SHAPE_TAGS:
            continue

        if _is_placeholder(el):
            idx = _get_ph_idx(el)
            # Skip footer/page-num
            if idx in _PH_SKIP:
                continue
            # Matching idx exists in dst -> skip, handled by _copy_placeholder_text
            if idx in dst_ph_idxs:
                continue
            # idx not in dst -> remove p:ph and transplant as free-floating textbox
            el_copy = copy.deepcopy(el)
            for ph_el in el_copy.findall(f".//{qn('p:ph')}"):
                ph_el.getparent().remove(ph_el)
            # If coordinates are layout-inherited only, explicitly embed resolved coordinates
            if idx in src_ph_map:
                _inject_xfrm(el_copy, src_ph_map[idx])
            # Transplant image rId for picture placeholder (p:pic)
            _fix_pic_rels(el_copy, src_slide, dst_slide)
            dst_tree.append(el_copy)
        else:
            _copy_shape_element(el, src_slide, dst_slide, dst_tree)


def _copy_placeholder_text(src_slide, dst_slide, dst_ph_idxs: set):
    """
    Copy placeholder text from src to dst for idx values that exist in dst.
    """
    for src_ph in src_slide.placeholders:
        idx = src_ph.placeholder_format.idx
        if idx in _PH_SKIP:
            continue
        if idx not in dst_ph_idxs:
            continue
        try:
            dst_ph = dst_slide.placeholders[idx]
        except KeyError:
            continue
        if src_ph.has_text_frame and dst_ph.has_text_frame:
            src_body = src_ph.text_frame._txBody
            dst_body = dst_ph.text_frame._txBody
            dst_body.getparent().replace(dst_body, copy.deepcopy(src_body))


# ── Auto theme detection ─────────────────────────────────────────────────────────────

def _auto_detect_source_theme(src_path: str) -> dict:
    """
    Auto-detect source theme by matching colors against all available themes, return tokens.
    Returns empty dict if no matches (continues without color replacement).
    """
    theme_files = [
        f for f in glob.glob(os.path.join(_THEMES_DIR, "*.json"))
        if not os.path.basename(f).startswith("_")
    ]
    if not theme_files:
        print("Warning: no theme files found, color replacement skipped")
        return {}

    from master_to_theme import _collect_colors_from_prs
    src_prs = Presentation(src_path)
    fill_c, text_c, _ = _collect_colors_from_prs(src_prs)
    src_colors = {c.lstrip("#").upper() for c in list(fill_c) + list(text_c)}

    best_name, best_tokens, best_score = None, {}, -1
    for tf in theme_files:
        with open(tf, encoding="utf-8") as f:
            data = json.load(f)
        tokens = data.get("tokens", {})
        score = sum(
            1 for v in tokens.values()
            if v.lstrip("#").upper() in src_colors
        )
        if score > best_score:
            best_score, best_name, best_tokens = score, data.get("name"), tokens

    if best_score > 0:
        print(f"Auto-detected source theme: {best_name} ({best_score} color matches)")
        return best_tokens

    print("Warning: could not auto-detect source theme, color replacement skipped")
    return {}


# ── Main processing ────────────────────────────────────────────────────────────────

def migrate(
    src_path: str,
    target_theme: str,
    source_theme: str | None = None,
    output_path: str | None = None,
) -> str:
    """
    Migrate the PPTX theme.

    Args:
        src_path:     Source PPTX
        target_theme: Target theme name (name field in themes/*.json)
        source_theme: Source theme name (auto-detected if omitted)
        output_path:  Output path (default: <input>_<target>.pptx)

    Returns:
        Output file path
    """
    if output_path is None:
        base = os.path.splitext(src_path)[0]
        output_path = f"{base}_{target_theme}.pptx"

    print(f"Source : {src_path}")
    print(f"Theme  : {source_theme or 'auto-detect'} → {target_theme}")
    print(f"Output : {output_path}")
    print()

    # ── Load themes ──
    dst_data     = _load_theme_data(target_theme)
    dst_tokens   = dst_data.get("tokens", {})
    dst_tpl_path = _resolve_template_path(dst_data)
    print(f"New template: {dst_tpl_path}")

    src_tokens = (
        _load_theme_data(source_theme).get("tokens", {})
        if source_theme
        else _auto_detect_source_theme(src_path)
    )

    # ── Presentations ──
    src_prs = Presentation(src_path)
    dst_prs = Presentation(dst_tpl_path)

    # Build schemeClr map from source theme
    src_theme_map = _build_src_theme_map(src_prs)

    # Collect actually used colors from slides (for fuzzy matching)
    actual_colors = _collect_slide_colors(src_prs)

    color_map = _build_color_map_fuzzy(src_tokens, dst_tokens, actual_colors)
    if color_map:
        print(f"\nColor map ({len(color_map)} replacements):")
        for old, new in sorted(color_map.items()):
            print(f"  #{old} → #{new}")
    else:
        print("\nNo color replacements (themes may use identical colors or detection failed)")
    print()

    # Clear existing slides from new template
    while len(dst_prs.slides):
        sld_id = dst_prs.slides._sldIdLst[0]
        r_id   = sld_id.get(qn("r:id"))
        dst_prs.part.drop_rel(r_id)
        del dst_prs.slides._sldIdLst[0]

    # Layout type map (new template)
    dst_layout_map = _build_dst_layout_map(dst_prs)

    # ── Copy slides ──
    print(f"Migrating {len(src_prs.slides)} slides...")
    for i, src_slide in enumerate(src_prs.slides, 1):
        src_layout_name = src_slide.slide_layout.name
        src_layout_type = _classify_layout_type(src_layout_name)
        new_layout      = dst_layout_map.get(src_layout_type, dst_layout_map["content"])
        print(f"  Slide {i:2d}: [{src_layout_name}] → {src_layout_type} → [{new_layout.name}]")

        dst_slide = dst_prs.slides.add_slide(new_layout)

        # Placeholder idx set of dst layout
        dst_ph_idxs = {
            ph.placeholder_format.idx
            for ph in dst_slide.placeholders
        }

        # For content slides, only copy title (idx=0) via placeholder.
        # idx=10 (body) / idx=11 (breadcrumb) may differ in position between templates,
        # so preserve original position as free-floating textbox.
        if src_layout_type == "content":
            dst_ph_idxs = dst_ph_idxs & {0}

        _copy_shapes(src_slide, dst_slide, dst_ph_idxs)
        _copy_placeholder_text(src_slide, dst_slide, dst_ph_idxs)

        # Remove empty placeholders (prevent "Click to add title" hint text from showing)
        _remove_empty_placeholders(dst_slide)

        # Resolve schemeClr to srgbClr, then apply color_map
        if src_theme_map:
            _resolve_scheme_colors(dst_slide._element, src_theme_map, color_map)

        # Replace srgbClr via color_map (both resolved schemeClr and original srgbClr)
        replaced = 0
        if color_map:
            replaced = _replace_colors(dst_slide._element, color_map)

        replaced_str = f", {replaced} color(s) replaced" if replaced else ""
        print(f"          → OK{replaced_str}")

    # ── Save ──
    dst_prs.save(output_path)
    print(f"\nSaved: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="PPTX theme migration tool — batch-replace colors and slide masters"
    )
    parser.add_argument("input",        help="Source PPTX file path")
    parser.add_argument("target_theme", help="Target theme name (e.g. fiori, accenture)")
    parser.add_argument("--from", dest="source_theme", metavar="THEME",
                        help="Source theme name (auto-detected if omitted)")
    parser.add_argument("--out",  metavar="FILE",
                        help="Output file path (default: input_<target>.pptx)")
    args = parser.parse_args()

    migrate(
        src_path=args.input,
        target_theme=args.target_theme,
        source_theme=args.source_theme,
        output_path=args.out,
    )


if __name__ == "__main__":
    main()
