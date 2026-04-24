"""
master_to_theme.py — Auto-generate color theme JSON from slide content

Scans actual slide elements (shape fills, text colors, line colors, table cells)
from a PowerPoint file, and maps them to 9 design tokens based on frequency,
luminance, and saturation. Content-based extraction ensures colors match actual usage.

Usage:
  # Command line
  python master_to_theme.py path/to/client.pptx "Client XYZ"

  # Python API
  from master_to_theme import extract_and_save
  extract_and_save("client.pptx", "Client XYZ")

Token mapping logic (content-based):
  primary      : Most frequent mid-luminance/high-saturation fill color (badges, headers, etc.)
  primary_dark : Dark variant of primary (primary at -50% luminance)
  primary_light: Light variant of primary (found in high luminance range or computed)
  background   : #FFFFFF (fixed)
  surface      : Very light chromatic fill (card backgrounds, etc.)
  text_heading : Darkest text color
  text_body    : Second darkest text color (or lightened text_heading)
  text_muted   : Mid-luminance text color (annotations, secondary labels)
  border       : Light gray lines/fills

Font extraction logic:
  font_heading : Extracted from <a:majorFont>. For system fonts, estimated from
                 title placeholder usage frequency
  font_body    : Extracted from <a:minorFont>. For system fonts, estimated from body run usage frequency
"""

import os
import sys
import json
import colorsys
from collections import Counter

_SCRIPTS_DIR = os.path.dirname(os.path.abspath(__file__))
_SKILL_DIR   = os.path.dirname(_SCRIPTS_DIR)
_THEMES_DIR  = os.path.join(_SKILL_DIR, "assets", "themes")


# ── Color utilities ──────────────────────────────────────────────────────────

def _parse_hex(hex_str: str) -> tuple:
    h = hex_str.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _to_hex(r: int, g: int, b: int) -> str:
    return f"#{r:02X}{g:02X}{b:02X}"


def _darken(hex_str: str, factor: float) -> str:
    """factor=0.5 -> darken to 50% brightness"""
    r, g, b = _parse_hex(hex_str)
    return _to_hex(int(r * factor), int(g * factor), int(b * factor))


def _lighten(hex_str: str, factor: float) -> str:
    """factor=0.85 -> lighten 85% toward white"""
    r, g, b = _parse_hex(hex_str)
    return _to_hex(
        int(r + (255 - r) * factor),
        int(g + (255 - g) * factor),
        int(b + (255 - b) * factor),
    )


def _luminance(hex_str: str) -> float:
    r, g, b = [c / 255 for c in _parse_hex(hex_str)]
    return 0.299 * r + 0.587 * g + 0.114 * b


def _saturation(hex_str: str) -> float:
    """HLS saturation (0-1)"""
    r, g, b = [c / 255 for c in _parse_hex(hex_str)]
    _, _, s = colorsys.rgb_to_hls(r, g, b)
    return s


def _hue(hex_str: str) -> float:
    """Hue (0-360 degrees)"""
    r, g, b = [c / 255 for c in _parse_hex(hex_str)]
    h, _, _ = colorsys.rgb_to_hls(r, g, b)
    return h * 360


def _hue_diff(h1: str, h2: str) -> float:
    """Hue difference between two colors (0-180 degrees)"""
    d = abs(_hue(h1) - _hue(h2))
    return min(d, 360 - d)


def _contrast_text(hex_bg: str) -> str:
    return "#FFFFFF" if _luminance(hex_bg) < 0.5 else "#000000"


def _is_near_white(hex_str: str, threshold: float = 0.92) -> bool:
    return _luminance(hex_str) >= threshold


def _is_near_black(hex_str: str, threshold: float = 0.12) -> bool:
    return _luminance(hex_str) <= threshold


# ── Font extraction ─────────────────────────────────────────────────────────────

# Office/OS default fonts (not brand fonts)
_SYSTEM_FONTS = frozenset({
    "Calibri", "Calibri Light", "Arial", "Arial Narrow", "Arial Black",
    "Times New Roman", "Helvetica", "Helvetica Neue", "Georgia", "Verdana",
    "Tahoma", "Trebuchet MS", "Symbol", "Wingdings", "Courier New",
    "Comic Sans MS", "Segoe UI", "Segoe UI Light", "Segoe UI Semibold",
    "Century Gothic", "Garamond", "Palatino Linotype", "Franklin Gothic Medium",
})
_OFFICE_PLACEHOLDERS = frozenset({"+mj-lt", "+mn-lt", "+mj-ea", "+mn-ea",
                                   "+mj-cs", "+mn-cs"})


def _is_system_font(name: str | None) -> bool:
    """Check if font is a system/Office default font"""
    return (name is None
            or name in _SYSTEM_FONTS
            or name in _OFFICE_PLACEHOLDERS
            or name.startswith("+"))


# ── Collect colors from slide content ───────────────────────────────────────────

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS   = {"a": _NS_A}


def _font_scheme_from_master(prs) -> tuple:
    """
    Read <a:fontScheme> from slide master theme XML.

    Returns:
        (heading_font, body_font) — None if not found or system font
    """
    for master in prs.slide_masters:
        for rel in master.part.rels.values():
            if "theme" not in rel.reltype.lower():
                continue
            from lxml import etree as _et
            tree = _et.fromstring(rel.target_part.blob)
            scheme = tree.find(".//a:fontScheme", _NS)
            if scheme is None:
                continue
            major = scheme.find("a:majorFont/a:latin", _NS)
            minor = scheme.find("a:minorFont/a:latin", _NS)
            heading = (major.get("typeface") or None) if major is not None else None
            body    = (minor.get("typeface") or None) if minor is not None else None
            return heading, body
    return None, None


def _count_fonts_in_slides(prs, max_slides: int = 15) -> tuple:
    """
    Count <a:latin typeface> usage frequency in title vs body across slides.

    Returns:
        (title_counter, body_counter)
    """
    title_c: Counter = Counter()
    body_c:  Counter = Counter()

    _TITLE_TYPES = {1, 3, 13}   # TITLE, CENTER_TITLE, VERTICAL_TITLE

    for slide in list(prs.slides)[:max_slides]:
        title_phs = {ph.placeholder_format.idx
                     for ph in slide.placeholders
                     if ph.placeholder_format.type.value in _TITLE_TYPES}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title = (hasattr(shape, "placeholder_format")
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.idx in title_phs)
            target = title_c if is_title else body_c

            for rpr in shape._element.iter(f"{{{_NS_A}}}rPr"):
                latin = rpr.find(f"{{{_NS_A}}}latin")
                if latin is not None:
                    tf = latin.get("typeface")
                    if tf:
                        target[tf] += 1

    return title_c, body_c


def _extract_fonts_from_prs(prs) -> dict:
    """
    Extract heading/body font names from a Presentation object.

    Strategy:
      1. Read from slide master <a:fontScheme> (intended brand font)
      2. For system fonts, fall back to title/body placeholder usage frequency
      3. Filter system fonts to surface brand fonts

    Returns:
        {"font_heading": str|None, "font_body": str|None}
    """
    heading, body = _font_scheme_from_master(prs)

    if _is_system_font(heading) or _is_system_font(body):
        title_c, body_c = _count_fonts_in_slides(prs)
        brand_title = [f for f, _ in title_c.most_common() if not _is_system_font(f)]
        brand_body  = [f for f, _ in body_c.most_common()  if not _is_system_font(f)]

        if _is_system_font(heading):
            heading = brand_title[0] if brand_title else (brand_body[0] if brand_body else None)
        if _is_system_font(body):
            body = brand_body[0] if brand_body else (brand_title[0] if brand_title else None)

    if _is_system_font(heading):
        heading = None
    if _is_system_font(body):
        body = None

    return {"font_heading": heading, "font_body": body}


def _build_theme_map(prs) -> dict:
    """
    Build a schemeClr name -> hex dict from slide master theme XML.
    WINDOW/WINDOWTEXT are OS-dependent, replaced with black/white.
    """
    theme_map: dict = {}
    for master in prs.slide_masters:
        for rel in master.part.rels.values():
            if "theme" not in rel.reltype.lower():
                continue
            from lxml import etree as _et
            tree = _et.fromstring(rel.target_part.blob)
            clr_el = tree.find(".//a:clrScheme", _NS)
            if clr_el is None:
                continue
            for child in clr_el:
                tag = child.tag.split("}")[1]
                for sub in child:
                    val = sub.get("val") or sub.get("lastClr", "")
                    if not val:
                        continue
                    if val.upper() in ("WINDOWTEXT", "WINDOW"):
                        theme_map[tag] = "#000000" if "dk" in tag.lower() or tag == "tx1" else "#FFFFFF"
                    else:
                        theme_map[tag] = "#" + val.upper()
            break  # use first master only
    # bg1/tx1 aliases
    if "bg1" not in theme_map:
        theme_map["bg1"] = theme_map.get("lt1", "#FFFFFF")
    if "tx1" not in theme_map:
        theme_map["tx1"] = theme_map.get("dk1", "#000000")
    return theme_map


def _resolve_color_el(el, theme_map: dict) -> str | None:
    """
    Resolve color element (srgbClr / schemeClr / sysClr) under <a:solidFill> to hex.
    Also computes schemeClr lumMod / lumOff / shade / tint modifiers.
    """
    for child in el:
        tag = child.tag.split("}")[1]

        if tag == "srgbClr":
            val = child.get("val", "")
            if len(val) == 6:
                return "#" + val.upper()

        elif tag == "sysClr":
            last = child.get("lastClr", "")
            if len(last) == 6:
                return "#" + last.upper()

        elif tag == "schemeClr":
            base = theme_map.get(child.get("val", ""))
            if not base:
                continue
            r0, g0, b0 = _parse_hex(base)
            h, l, s = colorsys.rgb_to_hls(r0 / 255, g0 / 255, b0 / 255)

            for mod in child:
                mod_tag = mod.tag.split("}")[1]
                v = int(mod.get("val", "100000")) / 100000
                if mod_tag == "lumMod":
                    l = max(0.0, min(1.0, l * v))
                elif mod_tag == "lumOff":
                    l = max(0.0, min(1.0, l + v))
                elif mod_tag == "shade":
                    l = max(0.0, min(1.0, l * v))
                elif mod_tag == "tint":
                    l = max(0.0, min(1.0, l + (1 - l) * v))

            r1, g1, b1 = colorsys.hls_to_rgb(h, l, s)
            return _to_hex(int(r1 * 255), int(g1 * 255), int(b1 * 255))

    return None


def _scan_xml_for_colors(element, theme_map: dict,
                          fill_c: Counter, text_c: Counter, line_c: Counter):
    """
    Scan entire XML tree for solidFill elements, classify by parent tag usage (fill/text/line).
    Uses namespace-agnostic .iter() to handle both DrawingML(a:) and
    PresentationML(p:) namespaces.
    """
    _FILL_PARENTS = {"spPr", "grpSpPr", "tcPr", "bgPr", "style"}
    _LINE_PARENTS = {"ln", "lnRef", "uFill", "lnStyleLst"}
    _TEXT_PARENTS = {"rPr", "endParaRPr", "defRPr", "pPr"}

    for sf in element.iter(f"{{{_NS_A}}}solidFill"):
        parent = sf.getparent()
        if parent is None:
            continue
        parent_tag = parent.tag.split("}")[1]
        c = _resolve_color_el(sf, theme_map)
        if not c:
            continue
        if parent_tag in _TEXT_PARENTS:
            text_c[c] += 1
        elif parent_tag in _LINE_PARENTS:
            line_c[c] += 1
        elif parent_tag in _FILL_PARENTS:
            fill_c[c] += 1
        # Ignore theme definitions like bgFillStyleLst, fmtScheme


def _collect_colors_from_prs(prs, max_slides: int = 10) -> tuple:
    """
    Scan all shapes, text, and tables at XML level across slides,
    resolving schemeClr (theme reference colors), and return
    (fill_counter, text_counter, line_counter).
    """
    fill_c: Counter = Counter()
    text_c: Counter = Counter()
    line_c: Counter = Counter()

    theme_map = _build_theme_map(prs)

    # Slide masters + layouts
    for master in prs.slide_masters:
        _scan_xml_for_colors(master._element, theme_map, fill_c, text_c, line_c)
        for layout in master.slide_layouts:
            _scan_xml_for_colors(layout._element, theme_map, fill_c, text_c, line_c)

    # Regular slides (up to max_slides)
    for slide in list(prs.slides)[:max_slides]:
        _scan_xml_for_colors(slide._element, theme_map, fill_c, text_c, line_c)

    return fill_c, text_c, line_c


# ── Color-to-token mapping ───────────────────────────────────────────────────

def _map_tokens_from_colors(fill_c: Counter, text_c: Counter, line_c: Counter) -> dict:
    """Auto-map to 9 tokens based on frequency, luminance, and saturation"""

    all_fills = fill_c.most_common()

    # ── primary candidates: mid-luminance (0.15-0.80) + high-saturation (>0.18) fills ──
    primary_candidates = [
        c for c, _ in all_fills
        if _saturation(c) > 0.18
        and 0.15 < _luminance(c) < 0.80
        and not _is_near_white(c)
    ]
    primary = primary_candidates[0] if primary_candidates else "#4472C4"

    # ── Colors in the same hue family as primary (hue diff < 45 deg) ──
    family = [(c, n) for c, n in all_fills if _hue_diff(c, primary) < 45]

    # primary_dark: same family, darker than primary
    darker = sorted(
        [c for c, _ in family if _luminance(c) < _luminance(primary) - 0.05],
        key=_luminance
    )
    primary_dark = darker[0] if darker else _darken(primary, 0.50)

    # primary_light: same family, high luminance (>0.80) with some saturation
    lighter = sorted(
        [c for c, _ in family
         if _luminance(c) > 0.80 and _saturation(c) > 0.05 and not _is_near_white(c)],
        key=_luminance, reverse=True
    )
    primary_light = lighter[0] if lighter else _lighten(primary, 0.82)

    # ── background: fixed #FFFFFF ──
    background = "#FFFFFF"

    # ── surface: very light chromatic fill (luminance > 0.90 and saturation > 0.03) ──
    surface_candidates = [
        c for c, _ in all_fills
        if _luminance(c) > 0.90 and _saturation(c) > 0.03 and not _is_near_white(c, 0.97)
    ]
    surface = surface_candidates[0] if surface_candidates else _lighten(primary, 0.93)

    # ── Text colors ──
    all_texts = text_c.most_common()

    # text_heading: darkest text (any saturation)
    dark_texts = sorted([c for c, _ in all_texts if _luminance(c) < 0.35], key=_luminance)
    text_heading = dark_texts[0] if dark_texts else "#000000"

    # text_body: dark (<0.40) and grayish (saturation <0.20) text
    gray_body = [
        c for c, _ in all_texts
        if _luminance(c) < 0.40 and _saturation(c) < 0.20
        and _luminance(c) > _luminance(text_heading) + 0.02
    ]
    text_body = gray_body[0] if gray_body else _lighten(text_heading, 0.20)

    # text_muted: mid-luminance (0.30-0.70) and grayish (saturation <0.20) text
    muted_candidates = [
        c for c, _ in all_texts
        if 0.30 < _luminance(c) < 0.70 and _saturation(c) < 0.20
    ]
    text_muted = muted_candidates[0] if muted_candidates else _lighten(text_heading, 0.50)

    # ── border: light gray (prefer lines, fallback to fills) ──
    border_from_lines = [
        c for c, _ in line_c.most_common()
        if 0.65 < _luminance(c) < 0.95 and _saturation(c) < 0.15
    ]
    border_from_fills = [
        c for c, _ in all_fills
        if 0.70 < _luminance(c) < 0.93 and _saturation(c) < 0.10
    ]
    border = (
        border_from_lines[0] if border_from_lines
        else border_from_fills[0] if border_from_fills
        else _lighten(text_heading, 0.80)
    )

    return {
        "primary":       primary,
        "primary_light": primary_light,
        "primary_dark":  primary_dark,
        "background":    background,
        "surface":       surface,
        "text_heading":  text_heading,
        "text_body":     text_body,
        "text_muted":    text_muted,
        "border":        border,
    }


# ── Main extraction API ────────────────────────────────────────────────────────────

def extract_colors_from_master(pptx_path: str) -> dict:
    """
    Collect colors from slide content (shapes, text, tables, slide masters)
    and map them to design tokens.

    Returns:
        {
          "scheme_name": str,
          "raw": {"fill": {...}, "text": {...}, "line": {...}},  # top 10 colors by frequency
          "tokens": {primary, primary_light, ...}
        }
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    fill_c, text_c, line_c = _collect_colors_from_prs(prs)
    tokens = _map_tokens_from_colors(fill_c, text_c, line_c)
    fonts  = _extract_fonts_from_prs(prs)

    # raw: record top 10 colors by frequency for debugging
    raw = {
        "fill_top10":  dict(fill_c.most_common(10)),
        "text_top10":  dict(text_c.most_common(10)),
        "line_top10":  dict(line_c.most_common(10)),
    }

    return {
        "scheme_name": os.path.basename(pptx_path),
        "raw":    raw,
        "tokens": tokens,
        "fonts":  fonts,
    }


_TOKEN_LABELS = {
    "primary":       "Main action / emphasis",
    "primary_light": "Background highlight",
    "primary_dark":  "Cover background / divider",
    "background":    "Slide background",
    "surface":       "Card / box background",
    "text_heading":  "Heading text",
    "text_body":     "Body text",
    "text_muted":    "Annotation / secondary text",
    "border":        "Border / divider line",
}
_TOKEN_ORDER = list(_TOKEN_LABELS.keys())


def _print_token_preview(result: dict, theme_name: str):
    """Print extracted tokens to terminal (no GUI)"""
    tokens = result["tokens"]
    fonts  = result.get("fonts", {})
    print(f"\n=== Theme color extraction result: {theme_name} ===")
    print(f"Source file: {result['scheme_name']}\n")
    print(f"{'Token':<16}  {'HEX':<9}  Purpose")
    print("-" * 52)
    for token in _TOKEN_ORDER:
        label = _TOKEN_LABELS.get(token, "")
        print(f"{token:<16}  {tokens[token]:<9}  {label}")
    print()
    print(f"{'Font':<16}  Name")
    print("-" * 52)
    print(f"{'font_heading':<16}  {fonts.get('font_heading') or '(not detected)'}")
    print(f"{'font_body':<16}  {fonts.get('font_body') or '(not detected)'}")
    print()


# ── Main API ────────────────────────────────────────────────────────────────

def extract_and_save(
    pptx_path: str,
    theme_name: str | None = None,
    preview: bool = True,
) -> str | None:
    """
    Extract colors from slide content and save as JSON.
    No GUI (Tkinter) is used.

    Args:
        pptx_path : Path to .pptx/.potx file to analyze
        theme_name: Theme name (auto-derived from filename if None)
        preview   : True = display extraction results in terminal before saving

    Returns:
        Path of saved JSON file
    """
    if theme_name is None:
        theme_name = os.path.splitext(os.path.basename(pptx_path))[0]

    result = extract_colors_from_master(pptx_path)

    if preview:
        _print_token_preview(result, theme_name)

    os.makedirs(_THEMES_DIR, exist_ok=True)
    safe_name = theme_name.lower().replace(" ", "-").replace("/", "-")
    out_path  = os.path.join(_THEMES_DIR, f"{safe_name}.json")

    # Copy slide master pptx to assets/masters/ (for template reuse)
    masters_dir = os.path.join(_SKILL_DIR, "assets", "masters")
    os.makedirs(masters_dir, exist_ok=True)
    master_dst  = os.path.join(masters_dir, f"{safe_name}.pptx")
    import shutil as _shutil
    _shutil.copy2(pptx_path, master_dst)
    template_rel = os.path.join("assets", "masters", f"{safe_name}.pptx")

    fonts = result.get("fonts", {})
    payload = {
        "name":        theme_name,
        "description": f"Auto-extracted from {os.path.basename(pptx_path)}",
        "template":    template_rel,
        "tokens":      result["tokens"],
        # Font name — null if undetectable (system fonts only)
        "font_heading": fonts.get("font_heading"),
        "font_body":    fonts.get("font_body"),
    }
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)

    return out_path


# ── Command-line execution ────────────────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python master_to_theme.py <path/to/template.pptx> [theme_name]")
        print("Example: python master_to_theme.py client.pptx \"Client XYZ\"")
        sys.exit(1)

    pptx = sys.argv[1]
    name = sys.argv[2] if len(sys.argv) >= 3 else None

    saved = extract_and_save(pptx, name)
    if saved:
        print(f"Saved: {saved}")
    else:
        print("Failed.")
        sys.exit(1)
